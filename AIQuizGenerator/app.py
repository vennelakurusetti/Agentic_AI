import os
import json
import streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()

st.set_page_config(page_title="Quizify", layout="centered", page_icon="⚡")

# ── Session defaults ──
if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = True

theme = "dark" if st.session_state.dark_mode else "light"

# ── Theme-aware CSS variables ──
CSS = f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600;700&family=Inter:wght@400;500;600&family=JetBrains+Mono:wght@400;500&display=swap');

:root {{
    --bg:          {"#0a0b12" if theme == "dark" else "#f0f2f8"};
    --bg-gradient: {"linear-gradient(135deg, #0a0b12 0%, #14162b 50%, #0a0b12 100%)" if theme == "dark" else "linear-gradient(135deg, #f0f2f8 0%, #ffffff 50%, #f0f2f8 100%)"};
    --surface:     {"rgba(18, 20, 32, 0.85)" if theme == "dark" else "rgba(255, 255, 255, 0.85)"};
    --surface-2:   {"rgba(28, 30, 48, 0.7)" if theme == "dark" else "rgba(240, 242, 248, 0.7)"};
    --surface-3:   {"rgba(38, 40, 64, 0.5)" if theme == "dark" else "rgba(230, 232, 240, 0.5)"};
    --border:      {"rgba(255,255,255,0.06)" if theme == "dark" else "rgba(0,0,0,0.05)"};
    --border-2:    {"rgba(255,255,255,0.10)" if theme == "dark" else "rgba(0,0,0,0.08)"};
    --text:        {"#eef0f6" if theme == "dark" else "#1a1a2e"};
    --text-sub:    {"#7b82a0" if theme == "dark" else "#6b7280"};
    --text-muted:  {"#4a5070" if theme == "dark" else "#9ca3af"};
    --amber:       {"#f59e0b" if theme == "dark" else "#d97706"};
    --amber-light: {"#fbbf24" if theme == "dark" else "#f59e0b"};
    --amber-dim:   {"rgba(245,158,11,0.10)" if theme == "dark" else "rgba(217,119,6,0.08)"};
    --amber-glow:  {"rgba(245,158,11,0.18)" if theme == "dark" else "rgba(217,119,6,0.15)"};
    --green:       {"#10b981" if theme == "dark" else "#059669"};
    --green-dim:   {"rgba(16,185,129,0.08)" if theme == "dark" else "rgba(5,150,105,0.06)"};
    --green-glow:  {"rgba(16,185,129,0.15)" if theme == "dark" else "rgba(5,150,105,0.12)"};
    --red:         {"#f43f5e" if theme == "dark" else "#dc2626"};
    --red-dim:     {"rgba(244,63,94,0.08)" if theme == "dark" else "rgba(220,38,38,0.06)"};
    --red-glow:    {"rgba(244,63,94,0.15)" if theme == "dark" else "rgba(220,38,38,0.12)"};
    --purple:      {"#8b5cf6" if theme == "dark" else "#7c3aed"};
    --purple-dim:  {"rgba(139,92,246,0.08)" if theme == "dark" else "rgba(124,58,237,0.06)"};
    --radius:      14px;
    --radius-lg:   18px;
    --radius-sm:   8px;
    --shadow-sm:   {"0 1px 3px rgba(0,0,0,0.3), 0 1px 2px rgba(0,0,0,0.2)" if theme == "dark" else "0 1px 3px rgba(0,0,0,0.04), 0 1px 2px rgba(0,0,0,0.06)"};
    --shadow-md:   {"0 8px 32px rgba(0,0,0,0.3), 0 2px 8px rgba(0,0,0,0.2)" if theme == "dark" else "0 8px 32px rgba(0,0,0,0.06), 0 2px 8px rgba(0,0,0,0.04)"};
    --shadow-lg:   {"0 20px 60px rgba(0,0,0,0.4), 0 0 40px rgba(245,158,11,0.08)" if theme == "dark" else "0 20px 60px rgba(0,0,0,0.08), 0 0 40px rgba(217,119,6,0.06)"};
    --glass-bg:    {"rgba(18, 20, 32, 0.6)" if theme == "dark" else "rgba(255, 255, 255, 0.6)"};
}}

/* ── Keyframes ── */
@keyframes fadeIn {{
    from {{ opacity: 0; transform: translateY(12px); }}
    to   {{ opacity: 1; transform: translateY(0); }}
}}
@keyframes fadeInUp {{
    from {{ opacity: 0; transform: translateY(24px); }}
    to   {{ opacity: 1; transform: translateY(0); }}
}}
@keyframes scaleIn {{
    from {{ opacity: 0; transform: scale(0.92); }}
    to   {{ opacity: 1; transform: scale(1); }}
}}
@keyframes pulse {{
    0%, 100% {{ transform: scale(1); opacity: 1; }}
    50%      {{ transform: scale(1.05); opacity: 0.85; }}
}}
@keyframes shimmer {{
    0%   {{ background-position: -200% 0; }}
    100% {{ background-position: 200% 0; }}
}}
@keyframes float {{
    0%, 100% {{ transform: translateY(0); }}
    50%      {{ transform: translateY(-8px); }}
}}
@keyframes spin {{ from {{ transform: rotate(0deg); }} to {{ transform: rotate(360deg); }} }}
@keyframes confetti-fall {{
    0%   {{ transform: translateY(-10px) rotate(0deg); opacity: 1; }}
    100% {{ transform: translateY(100px) rotate(720deg); opacity: 0; }}
}}
@keyframes slideInRight {{
    from {{ opacity: 0; transform: translateX(20px); }}
    to   {{ opacity: 1; transform: translateX(0); }}
}}
@keyframes slideInLeft {{
    from {{ opacity: 0; transform: translateX(-20px); }}
    to   {{ opacity: 1; transform: translateX(0); }}
}}

/* ── Reset ── */
*, *::before, *::after {{ box-sizing: border-box; }}

html, body,
[data-testid="stAppViewContainer"],
[data-testid="stApp"] {{
    background: var(--bg-gradient) !important;
    background-attachment: fixed !important;
    color: var(--text) !important;
    font-family: 'Inter', sans-serif;
}}

#MainMenu, footer, [data-testid="stHeader"] {{ display: none !important; }}

.block-container {{
    max-width: 680px !important;
    padding: 2rem 1.5rem 5rem !important;
}}

/* ── Decorative background orbs ── */
.main-deco {{
    position: fixed;
    top: -120px;
    right: -120px;
    width: 400px;
    height: 400px;
    border-radius: 50%;
    background: radial-gradient(circle, rgba(245,158,11,0.06) 0%, transparent 70%);
    pointer-events: none;
    z-index: 0;
    animation: float 8s ease-in-out infinite;
}}
.main-deco-2 {{
    position: fixed;
    bottom: -80px;
    left: -80px;
    width: 300px;
    height: 300px;
    border-radius: 50%;
    background: radial-gradient(circle, rgba(139,92,246,0.05) 0%, transparent 70%);
    pointer-events: none;
    z-index: 0;
    animation: float 10s ease-in-out infinite reverse;
}}

/* ── Type ── */
h1, h2, h3, h4 {{
    font-family: 'Space Grotesk', sans-serif !important;
    color: var(--text) !important;
    letter-spacing: -0.03em;
}}
h1 {{ 
    font-size: 2.4rem !important; 
    font-weight: 700 !important; 
    margin: 0 !important; 
    line-height: 1.1 !important;
    background: linear-gradient(135deg, var(--text) 0%, var(--amber-light) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}}
h2 {{ font-size: 1.15rem !important; font-weight: 600 !important; margin: 0 0 0.6rem !important; }}
p  {{ color: var(--text); line-height: 1.65; }}
[data-testid="stMarkdownContainer"] p {{ color: var(--text); }}

/* ── Divider ── */
hr {{
    border: none !important;
    border-top: 1px solid var(--border) !important;
    margin: 1.5rem 0 !important;
}}

/* ── Glass card base ── */
.glass-card {{
    background: var(--glass-bg);
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid var(--border);
    border-radius: var(--radius-lg);
    padding: 1.5rem;
    box-shadow: var(--shadow-md);
    position: relative;
    overflow: hidden;
}}
.glass-card::before {{
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 1px;
    background: linear-gradient(90deg, transparent, var(--amber-dim), transparent);
    opacity: 0.6;
}}

/* ── File uploader ── */
section[data-testid="stFileUploader"] {{
    background: var(--glass-bg) !important;
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1.5px dashed var(--border-2) !important;
    border-radius: var(--radius-lg) !important;
    padding: 2.5rem 1.5rem !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1) !important;
    cursor: pointer;
    position: relative;
    overflow: hidden;
}}
section[data-testid="stFileUploader"]::before {{
    content: '';
    position: absolute;
    inset: 0;
    background: linear-gradient(135deg, var(--amber-dim) 0%, transparent 50%, var(--purple-dim) 100%);
    opacity: 0;
    transition: opacity 0.3s ease;
}}
section[data-testid="stFileUploader"]:hover {{
    border-color: var(--amber) !important;
    transform: translateY(-2px);
    box-shadow: 0 0 0 4px var(--amber-dim), 0 12px 40px rgba(245,158,11,0.1) !important;
}}
section[data-testid="stFileUploader"]:hover::before {{
    opacity: 1;
}}
section[data-testid="stFileUploader"] button {{
    background: var(--amber-dim) !important;
    border: 1px solid rgba(245,158,11,0.3) !important;
    color: var(--amber) !important;
    border-radius: var(--radius-sm) !important;
    font-weight: 600 !important;
    font-size: 0.85rem !important;
    transition: all 0.2s ease !important;
    position: relative;
    z-index: 1;
}}
section[data-testid="stFileUploader"] button:hover {{
    background: rgba(245,158,11,0.18) !important;
    transform: scale(1.03);
}}
section[data-testid="stFileUploader"] [data-testid="stMarkdownContainer"] p {{
    color: var(--text-sub) !important;
    font-size: 0.85rem !important;
    position: relative;
    z-index: 1;
}}

/* ── Slider ── */
div[data-testid="stSlider"] {{ padding: 0.25rem 0; }}
div[data-testid="stSlider"] div[data-baseweb="slider"] > div:first-child {{
    background: var(--surface-2) !important;
    height: 4px !important;
    border-radius: 100px !important;
}}
div[data-testid="stSlider"] div[data-baseweb="slider"] div[role="slider"] {{
    background: var(--amber) !important;
    border: 2px solid var(--bg) !important;
    box-shadow: 0 0 0 3px var(--amber-glow), 0 2px 8px rgba(0,0,0,0.3) !important;
    width: 20px !important;
    height: 20px !important;
    transition: transform 0.2s ease !important;
}}
div[data-testid="stSlider"] div[data-baseweb="slider"] div[role="slider"]:hover {{
    transform: scale(1.15);
}}

/* ── Selectbox ── */
div[data-baseweb="select"] > div {{
    background: var(--glass-bg) !important;
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    border: 1px solid var(--border-2) !important;
    border-radius: 10px !important;
    color: var(--text) !important;
    font-size: 0.9rem !important;
    transition: all 0.2s ease !important;
}}
div[data-baseweb="select"] > div:hover,
div[data-baseweb="select"] > div:focus-within {{
    border-color: var(--amber) !important;
    box-shadow: 0 0 0 3px var(--amber-dim) !important;
}}
div[data-baseweb="select"] svg {{ color: var(--text-sub) !important; }}
div[data-baseweb="select"] input {{ color: var(--text) !important; }}
ul[data-testid="stSelectboxVirtualDropdown"],
div[data-baseweb="popover"] ul {{
    background: var(--surface) !important;
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid var(--border-2) !important;
    border-radius: 10px !important;
    box-shadow: var(--shadow-md);
}}
div[data-baseweb="popover"] li {{
    color: var(--text) !important;
    font-size: 0.875rem !important;
    transition: all 0.15s ease !important;
}}
div[data-baseweb="popover"] li:hover {{
    background: var(--amber-dim) !important;
    color: var(--amber) !important;
}}

/* ── Widget labels ── */
label[data-testid="stWidgetLabel"] p,
label[data-testid="stWidgetLabel"] {{
    color: var(--text-sub) !important;
    font-size: 0.72rem !important;
    font-weight: 600 !important;
    text-transform: uppercase !important;
    letter-spacing: 0.08em !important;
    font-family: 'Inter', sans-serif !important;
}}

/* ── Buttons ── */
div.stButton > button[kind="primary"],
div.stDownloadButton > button[kind="primary"] {{
    font-family: 'Space Grotesk', sans-serif !important;
    font-weight: 600 !important;
    font-size: 0.95rem !important;
    letter-spacing: -0.01em;
    background: linear-gradient(135deg, var(--amber) 0%, var(--amber-light) 100%) !important;
    color: #0c0e14 !important;
    border: none !important;
    border-radius: 10px !important;
    padding: 0.7rem 1.5rem !important;
    box-shadow: 0 4px 20px rgba(245,158,11,0.3), 0 0 40px rgba(245,158,11,0.1) !important;
    transition: all 0.2s cubic-bezier(0.4, 0, 0.2, 1) !important;
    position: relative;
    overflow: hidden;
}}
div.stButton > button[kind="primary"]::after {{
    content: '';
    position: absolute;
    inset: 0;
    background: linear-gradient(135deg, transparent 40%, rgba(255,255,255,0.15) 50%, transparent 60%);
    background-size: 200% 100%;
    animation: shimmer 3s ease-in-out infinite;
}}
div.stButton > button[kind="primary"]:hover {{
    transform: translateY(-2px);
    box-shadow: 0 8px 32px rgba(245,158,11,0.4), 0 0 60px var(--amber-glow) !important;
}}
div.stButton > button[kind="primary"]:active {{
    transform: translateY(0);
}}

div.stButton > button:not([kind="primary"]),
div.stDownloadButton > button:not([kind="primary"]) {{
    font-family: 'Inter', sans-serif !important;
    font-weight: 500 !important;
    font-size: 0.875rem !important;
    background: var(--glass-bg) !important;
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    color: var(--text) !important;
    border: 1px solid var(--border-2) !important;
    border-radius: 10px !important;
    padding: 0.55rem 1.1rem !important;
    transition: all 0.2s ease !important;
    box-shadow: var(--shadow-sm);
}}
div.stButton > button:not([kind="primary"]):hover {{
    border-color: var(--amber) !important;
    color: var(--amber) !important;
    background: var(--amber-dim) !important;
    transform: translateY(-1px);
    box-shadow: 0 4px 16px rgba(245,158,11,0.15);
}}
div.stButton > button:not([kind="primary"]):active {{
    transform: translateY(0);
}}

/* ── Radio cards ── */
div[data-testid="stRadio"] > div[role="radiogroup"] {{
    display: flex;
    flex-direction: column;
    gap: 0.55rem;
    animation: fadeIn 0.35s ease forwards;
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label {{
    background: var(--glass-bg) !important;
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    border: 1px solid var(--border-2) !important;
    border-radius: 12px !important;
    padding: 0.9rem 1.2rem !important;
    cursor: pointer;
    transition: all 0.2s cubic-bezier(0.4, 0, 0.2, 1) !important;
    font-size: 0.9rem !important;
    font-weight: 400;
    color: var(--text) !important;
    display: flex;
    align-items: center;
    gap: 0.75rem;
    box-shadow: var(--shadow-sm);
    position: relative;
    overflow: hidden;
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label::before {{
    content: '';
    position: absolute;
    inset: 0;
    background: linear-gradient(135deg, var(--amber-dim) 0%, transparent 100%);
    opacity: 0;
    transition: opacity 0.2s ease;
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label:hover {{
    border-color: var(--amber) !important;
    transform: translateX(4px);
    box-shadow: 0 4px 16px rgba(245,158,11,0.12);
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label:hover::before {{
    opacity: 1;
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label:has(input:checked) {{
    border-color: var(--amber) !important;
    background: var(--amber-dim) !important;
    box-shadow: 0 0 0 3px var(--amber-glow);
    font-weight: 500 !important;
    transform: translateX(4px);
}}
div[data-testid="stRadio"] > div[role="radiogroup"] label:has(input:checked)::before {{
    opacity: 1;
}}
div[data-testid="stRadio"] > div[role="radiogroup"] input[type="radio"] {{
    accent-color: var(--amber) !important;
    transform: scale(1.2);
    flex-shrink: 0;
    position: relative;
    z-index: 1;
}}

/* ── Progress bar ── */
div[data-testid="stProgress"] > div {{
    height: 6px !important;
    border-radius: 100px !important;
    background: var(--surface-2) !important;
    overflow: hidden;
    box-shadow: inset 0 1px 3px rgba(0,0,0,0.1);
}}
div[data-testid="stProgress"] > div > div {{
    background: linear-gradient(90deg, var(--amber) 0%, var(--amber-light) 100%) !important;
    border-radius: 100px;
    transition: width 0.6s cubic-bezier(0.4, 0, 0.2, 1);
    box-shadow: 0 0 12px var(--amber-glow);
    position: relative;
}}
div[data-testid="stProgress"] > div > div::after {{
    content: '';
    position: absolute;
    right: 0;
    top: 0;
    bottom: 0;
    width: 20px;
    background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3));
    border-radius: 100px;
}}

/* ── Alerts / info ── */
div.stAlert {{
    background: var(--surface) !important;
    border-radius: 12px !important;
    border: 1px solid var(--border-2) !important;
    color: var(--text) !important;
    font-size: 0.85rem !important;
    padding: 0.75rem 1rem !important;
    animation: fadeIn 0.3s ease;
}}
div[data-testid="stNotification"] {{
    background: var(--surface) !important;
    border-color: var(--border-2) !important;
    border-radius: 10px !important;
}}

/* ── Bordered container (question cards) ── */
div[data-testid="stVerticalBlock"] > div.stBorder {{
    background: var(--glass-bg) !important;
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid var(--border) !important;
    border-radius: var(--radius-lg) !important;
    padding: 1.5rem 1.5rem 1.25rem !important;
    margin: 0.5rem 0 !important;
    position: relative;
    overflow: hidden;
    box-shadow: var(--shadow-md);
    animation: scaleIn 0.4s cubic-bezier(0.4, 0, 0.2, 1) forwards;
}}
div[data-testid="stVerticalBlock"] > div.stBorder::before {{
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 1px;
    background: linear-gradient(90deg, transparent, var(--amber-dim), var(--amber-dim), transparent);
    opacity: 0.5;
}}

/* ── Question number ghost ── */
.q-num {{
    position: absolute;
    right: 1rem;
    top: -0.6rem;
    font-family: 'Space Grotesk', sans-serif;
    font-size: 4.5rem;
    font-weight: 700;
    color: {"rgba(245,158,11,0.06)" if theme == "dark" else "rgba(217,119,6,0.08)"};
    line-height: 1;
    user-select: none;
    pointer-events: none;
    letter-spacing: -0.04em;
}}

/* ── Question text ── */
.question-text {{
    font-family: 'Space Grotesk', sans-serif;
    font-size: 1.1rem;
    font-weight: 600;
    letter-spacing: -0.02em;
    margin: 0 0 1rem 0;
    line-height: 1.5;
    color: var(--text);
}}

/* ── Result card overlays ── */
.card-correct {{
    outline: 1.5px solid {"rgba(16,185,129,0.35)" if theme == "dark" else "rgba(5,150,105,0.3)"} !important;
    background: {"rgba(16,185,129,0.06)" if theme == "dark" else "rgba(5,150,105,0.04)"} !important;
    box-shadow: 0 0 24px var(--green-glow) !important;
}}
.card-wrong {{
    outline: 1.5px solid {"rgba(244,63,94,0.35)" if theme == "dark" else "rgba(220,38,38,0.3)"} !important;
    background: {"rgba(244,63,94,0.06)" if theme == "dark" else "rgba(220,38,38,0.04)"} !important;
    box-shadow: 0 0 24px var(--red-glow) !important;
}}

/* ── Score display ── */
.score-wrap {{
    text-align: center;
    padding: 2rem 0 1rem;
    animation: scaleIn 0.6s cubic-bezier(0.4, 0, 0.2, 1) forwards;
}}
.score-num {{
    font-family: 'Space Grotesk', sans-serif;
    font-size: 5rem;
    font-weight: 700;
    letter-spacing: -0.05em;
    background: linear-gradient(135deg, var(--amber) 0%, var(--amber-light) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    line-height: 1;
    display: inline-block;
}}
.score-denom {{
    font-family: 'Space Grotesk', sans-serif;
    font-size: 1.8rem;
    font-weight: 600;
    color: var(--text-muted);
}}
.score-label {{
    font-size: 0.9rem;
    color: var(--text-sub);
    margin-top: 0.5rem;
    font-weight: 500;
}}
.score-emoji {{
    font-size: 3rem;
    margin-bottom: 0.5rem;
    display: block;
    animation: float 3s ease-in-out infinite;
}}

/* ── Stat pills ── */
.pill-row {{
    display: flex;
    gap: 0.5rem;
    flex-wrap: wrap;
    margin: 0.75rem 0;
    justify-content: center;
}}
.pill {{
    background: var(--glass-bg);
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    border: 1px solid var(--border);
    border-radius: 100px;
    padding: 0.3rem 0.85rem;
    font-size: 0.75rem;
    color: var(--text-sub);
    font-weight: 500;
    display: inline-flex;
    align-items: center;
    gap: 0.35rem;
    transition: all 0.2s ease;
    animation: fadeIn 0.4s ease forwards;
    opacity: 0;
    animation-fill-mode: forwards;
}}
.pill:nth-child(1) {{ animation-delay: 0.1s; }}
.pill:nth-child(2) {{ animation-delay: 0.2s; }}
.pill:nth-child(3) {{ animation-delay: 0.3s; }}
.pill:nth-child(4) {{ animation-delay: 0.4s; }}
.pill:nth-child(5) {{ animation-delay: 0.5s; }}
.pill:nth-child(6) {{ animation-delay: 0.6s; }}
.pill:hover {{
    border-color: var(--amber-dim);
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(245,158,11,0.1);
}}
.pill b {{ color: var(--text); font-weight: 600; }}

/* ── Text preview ── */
div[data-testid="stText"],
div[data-testid="stCode"] {{
    background: var(--surface-2) !important;
    border: 1px solid var(--border) !important;
    border-radius: 10px !important;
    color: var(--text-sub) !important;
    font-size: 0.78rem !important;
    font-family: 'JetBrains Mono', monospace !important;
    line-height: 1.6 !important;
}}

/* ── Expander ── */
details summary {{
    color: var(--text-sub) !important;
    font-size: 0.82rem !important;
    font-weight: 500 !important;
}}
div[data-testid="stExpander"] {{
    background: var(--surface) !important;
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    border: 1px solid var(--border) !important;
    border-radius: 12px !important;
    overflow: hidden;
    transition: all 0.2s ease !important;
}}

/* ── Spinner ── */
div.stSpinner > div {{ 
    color: var(--text-sub) !important; 
    animation: spin 1s linear infinite !important;
}}

/* ── Nav meta ── */
.nav-meta {{
    font-size: 0.78rem;
    color: var(--text-muted);
    text-align: right;
    margin-bottom: 0.35rem;
    font-variant-numeric: tabular-nums;
    font-weight: 500;
    letter-spacing: 0.02em;
}}

/* ── Columns ── */
div.row-widget.stColumns {{ gap: 0.7rem; }}

/* ── Toggle ── */
[data-testid="stToggle"] {{
    display: flex !important;
    justify-content: flex-end;
}}
[data-testid="stToggle"] label {{
    gap: 0.4rem !important;
    font-size: 0.85rem !important;
    color: var(--text-sub) !important;
}}
[data-testid="stToggle"] label span:first-child {{
    font-size: 1.1rem !important;
}}
/* toggle track */
[data-testid="stToggle"] div[data-baseweb="toggle"] {{
    background: var(--surface-2) !important;
    border: 1px solid var(--border-2) !important;
    transition: all 0.2s ease !important;
}}
[data-testid="stToggle"] div[data-baseweb="toggle"][aria-checked="true"] {{
    background: linear-gradient(135deg, var(--amber) 0%, var(--amber-light) 100%) !important;
    border-color: var(--amber) !important;
    box-shadow: 0 0 12px var(--amber-glow);
}}

/* ── Subtitle animation ── */
.subtitle {{
    color: var(--text-sub);
    font-size: 0.95rem;
    margin: 0 0 1.5rem 0;
    animation: fadeIn 0.6s ease 0.2s forwards;
    opacity: 0;
}}

/* ── Hint badge ── */
.hint-badge {{
    font-size: 0.75rem;
    color: var(--text-muted);
    padding: 0.55rem 0;
    display: flex;
    align-items: center;
    gap: 0.4rem;
}}
.hint-badge::before {{
    content: '💡';
}}

/* ── Result item icon animation ── */
.result-icon {{
    display: inline-block;
    animation: scaleIn 0.3s ease forwards;
}}

/* ── Confetti container ── */
.confetti-container {{
    position: fixed;
    top: 0;
    left: 0;
    width: 100%;
    height: 100%;
    pointer-events: none;
    overflow: hidden;
    z-index: 9999;
}}
.confetti {{
    position: absolute;
    width: 10px;
    height: 10px;
    border-radius: 2px;
    animation: confetti-fall 3s ease-in forwards;
}}
.confetti:nth-child(1) {{ left: 10%; background: var(--amber); animation-delay: 0s; }}
.confetti:nth-child(2) {{ left: 25%; background: var(--green); animation-delay: 0.15s; }}
.confetti:nth-child(3) {{ left: 40%; background: var(--purple); animation-delay: 0.3s; }}
.confetti:nth-child(4) {{ left: 55%; background: var(--red); animation-delay: 0.45s; }}
.confetti:nth-child(5) {{ left: 70%; background: var(--amber-light); animation-delay: 0.6s; }}
.confetti:nth-child(6) {{ left: 85%; background: var(--green); animation-delay: 0.75s; }}
.confetti:nth-child(7) {{ left: 50%; background: var(--purple); animation-delay: 0.9s; }}
.confetti:nth-child(8) {{ left: 15%; background: var(--red); animation-delay: 1.05s; }}
.confetti:nth-child(9) {{ left: 65%; background: var(--amber); animation-delay: 1.2s; }}
.confetti:nth-child(10) {{ left: 35%; background: var(--green); animation-delay: 1.35s; }}
.confetti:nth-child(11) {{ left: 75%; background: var(--purple); animation-delay: 1.5s; }}
.confetti:nth-child(12) {{ left: 5%; background: var(--amber-light); animation-delay: 1.65s; }}

/* ── Info callout ── */
.info-callout {{
    background: {"rgba(139,92,246,0.06)" if theme == "dark" else "rgba(124,58,237,0.04)"};
    border: 1px solid {"rgba(139,92,246,0.12)" if theme == "dark" else "rgba(124,58,237,0.1)"};
    border-radius: 10px;
    padding: 0.7rem 1rem;
    font-size: 0.82rem;
    color: var(--text-sub);
    line-height: 1.5;
    margin-top: 0.5rem;
    animation: fadeIn 0.3s ease forwards;
}}

@media (max-width: 520px) {{
    .block-container {{ padding: 1.5rem 1rem 4rem !important; }}
    h1 {{ font-size: 1.8rem !important; }}
    .score-num {{ font-size: 3.5rem !important; }}
    .score-emoji {{ font-size: 2.2rem; }}
    .question-text {{ font-size: 1rem !important; }}
    div[data-testid="stVerticalBlock"] > div.stBorder {{ padding: 1.25rem 1rem 1rem !important; }}
}}

/* ── Staggered fade-in for result cards ── */
.result-card {{
    animation: fadeInUp 0.5s ease forwards;
    opacity: 0;
}}
.result-card:nth-child(1) {{ animation-delay: 0.05s; }}
.result-card:nth-child(2) {{ animation-delay: 0.1s; }}
.result-card:nth-child(3) {{ animation-delay: 0.15s; }}
.result-card:nth-child(4) {{ animation-delay: 0.2s; }}
.result-card:nth-child(5) {{ animation-delay: 0.25s; }}
.result-card:nth-child(6) {{ animation-delay: 0.3s; }}
.result-card:nth-child(7) {{ animation-delay: 0.35s; }}
.result-card:nth-child(8) {{ animation-delay: 0.4s; }}
.result-card:nth-child(9) {{ animation-delay: 0.45s; }}
.result-card:nth-child(10) {{ animation-delay: 0.5s; }}
.result-card:nth-child(11) {{ animation-delay: 0.55s; }}
.result-card:nth-child(12) {{ animation-delay: 0.6s; }}
.result-card:nth-child(13) {{ animation-delay: 0.65s; }}
.result-card:nth-child(14) {{ animation-delay: 0.7s; }}
.result-card:nth-child(15) {{ animation-delay: 0.75s; }}
.result-card:nth-child(16) {{ animation-delay: 0.8s; }}
.result-card:nth-child(17) {{ animation-delay: 0.85s; }}
.result-card:nth-child(18) {{ animation-delay: 0.9s; }}
.result-card:nth-child(19) {{ animation-delay: 0.95s; }}
.result-card:nth-child(20) {{ animation-delay: 1.0s; }}
.result-card:nth-child(21) {{ animation-delay: 1.05s; }}
.result-card:nth-child(22) {{ animation-delay: 1.1s; }}
.result-card:nth-child(23) {{ animation-delay: 1.15s; }}
.result-card:nth-child(24) {{ animation-delay: 1.2s; }}
.result-card:nth-child(25) {{ animation-delay: 1.25s; }}
.result-card:nth-child(26) {{ animation-delay: 1.3s; }}
.result-card:nth-child(27) {{ animation-delay: 1.35s; }}
.result-card:nth-child(28) {{ animation-delay: 1.4s; }}
.result-card:nth-child(29) {{ animation-delay: 1.45s; }}
.result-card:nth-child(30) {{ animation-delay: 1.5s; }}
</style>
"""

st.markdown(CSS, unsafe_allow_html=True)


# ── Helpers ──
def extract_text_from_ppt(prs) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
    return "\n".join(parts)


def pill(icon, label, value):
    return f'<div class="pill">{icon} {label} <b>{value}</b></div>'


# ══════════════════════════════════════════
#  DECORATIVE BACKGROUND
# ══════════════════════════════════════════
st.markdown(
    '<div class="main-deco"></div><div class="main-deco-2"></div>',
    unsafe_allow_html=True,
)


# ══════════════════════════════════════════
#  SETUP SCREEN
# ══════════════════════════════════════════
if "quiz" not in st.session_state:

    # ── Header row with toggle ──
    hcol1, hcol2 = st.columns([5, 1])
    with hcol1:
        st.markdown("""
        <div style="display:flex;align-items:center;gap:0.7rem;">
            <span style="font-size:1.8rem;animation:float 3s ease-in-out infinite;display:inline-block;">⚡</span>
            <h1>Quizify</h1>
        </div>
        """, unsafe_allow_html=True)
    with hcol2:
        dark = st.toggle("🌙", value=st.session_state.dark_mode, key="theme_toggle")
        if dark != st.session_state.dark_mode:
            st.session_state.dark_mode = dark
            st.rerun()

    st.markdown("""
    <p class="subtitle">
        Transform your PowerPoint presentations into interactive quizzes in seconds.
    </p>
    """, unsafe_allow_html=True)

    st.markdown("<hr>", unsafe_allow_html=True)

    # ── Upload ──
    uploaded_file = st.file_uploader(
        "Upload your PowerPoint file to get started",
        type=["pptx"],
        label_visibility="collapsed"
    )

    st.markdown("<div style='height:1.25rem'></div>", unsafe_allow_html=True)

    # ── Controls in glass card ──
    st.markdown('<div class="glass-card">', unsafe_allow_html=True)
    col_a, col_b = st.columns(2)
    with col_a:
        num_questions = st.slider("Questions", min_value=5, max_value=30, value=10)
    with col_b:
        difficulty = st.selectbox("Difficulty", ["Simple", "Medium", "Complex"])
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("<div style='height:0.5rem'></div>", unsafe_allow_html=True)

    # ── File info + generate ──
    if uploaded_file is not None:
        prs = Presentation(uploaded_file)
        slide_count = len(prs.slides)
        full_text = extract_text_from_ppt(prs)
        word_count = len(full_text.split())

        st.markdown(
            f'<div class="pill-row">'
            f'{pill("📎", "", uploaded_file.name)}'
            f'{pill("🗂", "", f"{slide_count} slides")}'
            f'{pill("📝", "", f"{word_count:,} words")}'
            f'</div>',
            unsafe_allow_html=True,
        )

        with st.expander("Preview extracted text"):
            st.text(full_text[:700] + ("…" if len(full_text) > 700 else ""))

        st.markdown("<div style='height:0.5rem'></div>", unsafe_allow_html=True)

        if st.button("✨ Generate Quiz", use_container_width=True, type="primary"):
            with st.spinner("🧠 Analyzing your content and building a quiz…"):
                client = OpenAI(
                    base_url="https://openrouter.ai/api/v1",
                    api_key=os.getenv("OPENROUTER_API_KEY"),
                )
                prompt = (
                    f"Based on the following text, generate {num_questions} multiple-choice questions "
                    f"at {difficulty} difficulty level. "
                    "Return ONLY a JSON array of objects, each with keys: "
                    '"question", "options" (array of 4 strings), "answer", and "explanation". '
                    "Do not include any other text.\n\n"
                    f"Text:\n{full_text}"
                )
                resp = client.chat.completions.create(
                    model="openai/gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                )
                raw = resp.choices[0].message.content.strip()
                qjson = raw
                if qjson.startswith("```"):
                    qjson = qjson.split("\n", 1)[1]
                    if qjson.endswith("```"):
                        qjson = qjson[:-3]
                    qjson = qjson.strip()
                try:
                    st.session_state.quiz = json.loads(qjson)
                except json.JSONDecodeError:
                    st.error("Couldn't parse the quiz. Please try again.")
                    st.code(raw)
                    st.stop()

                st.session_state.q_idx = 0
                st.session_state.answers = {}
                st.session_state.submitted = False
                st.session_state.hints_used = {}
                st.session_state.total_penalty = 0.0
                st.rerun()

    else:
        st.button("✨ Generate Quiz", use_container_width=True, type="primary", disabled=True)


# ══════════════════════════════════════════
#  QUIZ SCREEN
# ══════════════════════════════════════════
elif not st.session_state.get("submitted"):
    quiz  = st.session_state.quiz
    idx   = st.session_state.q_idx
    total = len(quiz)
    item  = quiz[idx]

    # Progress
    st.progress(idx / total)
    st.markdown(
        f'<p class="nav-meta" style="animation:fadeIn 0.3s ease forwards;">Question {idx + 1} of {total}</p>',
        unsafe_allow_html=True,
    )

    # Hint logic: which option to eliminate
    eliminated = None
    if st.session_state.hints_used.get(idx):
        correct = item["answer"]
        for opt in item["options"]:
            if opt != correct:
                eliminated = opt
                break

    display_options = [o for o in item["options"] if o != eliminated] if eliminated else item["options"]

    # Question card
    with st.container(border=True):
        st.markdown(f'<div class="q-num">{idx + 1:02d}</div>', unsafe_allow_html=True)
        st.markdown(
            f'<p class="question-text">{item["question"]}</p>',
            unsafe_allow_html=True,
        )

        saved = st.session_state.answers.get(idx)
        saved_idx = display_options.index(saved) if saved in display_options else None

        choice = st.radio(
            "options",
            display_options,
            key=f"r_{idx}",
            index=saved_idx,
            label_visibility="collapsed",
        )

    # Hint button row
    hint_col1, hint_col2 = st.columns([1.2, 5])
    with hint_col1:
        hint_used = st.session_state.hints_used.get(idx, False)
        if not hint_used:
            if st.button("💡 Hint", use_container_width=True):
                st.session_state.hints_used[idx] = True
                st.session_state.total_penalty = st.session_state.get("total_penalty", 0.0) + 0.5
                st.rerun()
        else:
            st.markdown(
                f'<div class="hint-badge">Hint used (−0.5) — one wrong option removed</div>',
                unsafe_allow_html=True,
            )

    # Nav buttons
    is_last = idx == total - 1
    col1, col2, spacer, col3 = st.columns([1.2, 1.4, 2.5, 1.1])
    with col1:
        if idx > 0 and st.button("← Back", use_container_width=True):
            if choice:
                st.session_state.answers[idx] = choice
            st.session_state.q_idx -= 1
            st.rerun()
    with col2:
        label = "✓ Submit Quiz" if is_last else "Next →"
        if st.button(label, use_container_width=True, type="primary" if is_last else "secondary"):
            if choice:
                st.session_state.answers[idx] = choice
                if is_last:
                    st.session_state.submitted = True
                else:
                    st.session_state.q_idx += 1
                st.rerun()
            else:
                st.warning("Please select an answer to continue.")
    with col3:
        if st.button("✕ Quit", use_container_width=True):
            for k in ["quiz", "q_idx", "answers", "submitted", "hints_used", "total_penalty"]:
                st.session_state.pop(k, None)
            st.rerun()


# ══════════════════════════════════════════
#  RESULTS SCREEN
# ══════════════════════════════════════════
else:
    quiz  = st.session_state.quiz
    total = len(quiz)
    score = sum(
        1 for i, item in enumerate(quiz)
        if st.session_state.answers.get(i) == item["answer"]
    )
    pct   = (score / total) * 100
    wrong = total - score

    emoji = "🏆" if pct >= 80 else "👍" if pct >= 50 else "📚"
    msg   = "Excellent work! You've mastered this topic." if pct >= 80 else "Good effort! Keep practicing to improve." if pct >= 50 else "Keep studying! Review the material and try again."

    # Confetti for high scores
    if pct >= 80:
        st.markdown("""
        <div class="confetti-container">
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
            <div class="confetti"></div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown(
        f"""
        <div class="score-wrap">
            <span class="score-emoji">{emoji}</span>
            <div>
                <span class="score-num">{score}</span>
                <span class="score-denom"> / {total}</span>
            </div>
            <p class="score-label">{pct:.0f}% · {msg}</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    penalty = st.session_state.get("total_penalty", 0.0)
    hints_count = len(st.session_state.get("hints_used", {}))
    net_score = score - penalty

    st.markdown(
        f'<div class="pill-row">'
        f'{pill("✅", "Correct", score)}'
        f'{pill("❌", "Wrong", wrong)}'
        f'{pill("📊", "Score", f"{pct:.0f}%")}'
        + (f'{pill("💡", "Hints", hints_count)}'
           f'{pill("⚠️", "Penalty", f"-{penalty:.1f}")}'
           f'{pill("🎯", "Net", f"{net_score:.1f}")}' if hints_count else "")
        + f'</div>',
        unsafe_allow_html=True,
    )

    st.markdown("<div style='height:0.75rem'></div>", unsafe_allow_html=True)
    st.progress(score / total)
    st.markdown("<div style='height:1.25rem'></div>", unsafe_allow_html=True)

    for i, item in enumerate(quiz):
        user  = st.session_state.answers.get(i, "—")
        ok    = user == item["answer"]
        icon  = "✅" if ok else "❌"

        with st.container(border=True):
            st.markdown(
                f'<div class="result-card">'
                f'<div class="q-num">{i + 1:02d}</div>'
                f'<p class="question-text" style="font-size:0.95rem;">'
                f'<span class="result-icon">{icon}</span> {item["question"]}</p>'
                f'<p style="font-size:0.82rem;color:var(--text-sub);margin:0 0 0.2rem;">'
                f'Your answer: <span style="color:var(--text);font-weight:500;">{user}</span></p>',
                unsafe_allow_html=True,
            )
            if not ok:
                st.markdown(
                    f'<p style="font-size:0.82rem;color:var(--green);margin:0 0 0.5rem;">'
                    f'Correct answer: <b>{item["answer"]}</b></p>',
                    unsafe_allow_html=True,
                )
            st.markdown(
                f'<div class="info-callout">💡 {item["explanation"]}</div>'
                f'</div>',
                unsafe_allow_html=True,
            )

    st.markdown("<div style='height:1.5rem'></div>", unsafe_allow_html=True)

    if st.button("↺  Take Another Quiz", use_container_width=True, type="primary"):
        for k in ["quiz", "q_idx", "answers", "submitted", "hints_used", "total_penalty"]:
            st.session_state.pop(k, None)
        st.rerun()